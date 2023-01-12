import collections
import collections.abc
import re

from pptx import Presentation
import pprint



'''
def main():
    prints text of slides that include a 'Lesson' word [[text in slide n], [text in slide m], ...]

def extract_text_from_ppt(ppt_path):
    extracts all text in a ppt file. returns [[text in slide 1], [text in slide 2], ...]

def filter_lesson_lists(slides_text):
    filters lists that a 'Lesson' word. returns [[text in slide n], [text in slide m], ...]

'''


def main():
    ppt_path = "Chatbot AGP 3 Presentation.pptx"
    data = filter_lesson_lists(extract_text_from_ppt(ppt_path))
    #pprint.pprint(extract_text_from_ppt(ppt_path))
    pprint.pprint(data)


# need to find a way to get data from inside table???
def extract_text_from_ppt(ppt_path):
    # Open the PowerPoint file using python-pptx
    presentation = Presentation(ppt_path)

    # Initialize an empty list to store the text
    slides_text = []

    # Loop through all slides in the presentation
    for slide in presentation.slides:
        # Loop through all shapes in the slide
        text_list = " "
        for shape in slide.shapes:
            # Check if the shape is a text box
            if shape.has_text_frame:
                # Extract the text from the text box and add it to the list
                text = shape.text
                # remove appendix because they will also have the keywords the algorithm will detect
                if 'appendix' in text.lower():
                    break
                elif 'sensitivity' not in text.lower():
                    text_list = text_list + str(text)
        slides_text.append(text_list)
    return slides_text


# Takes Action Items and lessons learned from the read slides text and appends them onto a new list
def filter_lesson_lists(slides_text):
    # new_slides_text slices the title page out
    new_slides_text = slides_text[1:]
    filtered_slides_text = []
    for text in new_slides_text:
        if 'lessons learned' in text.lower():
            key_word = re.search(r'lessons learned', text, re.IGNORECASE)
            filtered_slides_text.append(text)
        #    filtered_slides_text.append(text[key_word.start():])
        # elif 'action item' in text.lower():
        #     key_word = re.search(r'action item', text, re.IGNORECASE)
        #     filtered_slides_text.append("action item : \n")
        #     filtered_slides_text.append(text[key_word.end():])
        # elif 'decision criteria' in text.lower():
        #     key_word = re.search(r'Decision Criteria', text, re.IGNORECASE)
        #     filtered_slides_text.append("Descision Criteria : \n")
        #     filtered_slides_text.append(text[key_word.end():])

    newFile = open("temp.txt", "w")
    for text in filtered_slides_text:
        newFile.write(text)
    newFile.close()

    return filtered_slides_text


if __name__ == "__main__":
    main()
