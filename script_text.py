import collections
import collections.abc
import torch as pt
from pptx import Presentation
import pprint
import os
import re
import os.path
from transformers import pipeline
from docx import Document


def path_adder(dir_path):
    # list to store files
    res = []

    # Iterate directory
    for path in os.listdir(dir_path):
        # check if current path is a file
        if os.path.isfile(os.path.join(dir_path, path)):
            res.append(path)
    return res

def extract_text_from_ppt(presentation):
    prs = presentation
    # list_of slides will have all the text extracted from slides
    list_of_slides = []
    for slide in prs.slides:
        # list_of_elements stores all the sentences of one shape in slide
        list_of_elements = []
        # Go through all the shapes in the slides
        for shape in slide.shapes:
            # Check if the shape is in table form
            if shape.has_table:
                line = ''
                for cell in shape.table.iter_cells():
                    line += cell.text+'\n'
                list_of_elements.append(line)
            # Check if Shape is in textframe
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = ''
                    for run in paragraph.runs:
                        if "sensitivity" not in run.text.lower():
                            line += run.text
                    list_of_elements.append(line)

        # remove all the extra characters ' ' and ''
        new_list = [x for x in list_of_elements if x != '' and x != ' ']
        list_of_slides.append(new_list)
    return list_of_slides

# Takes Action Items and lessons learned from the read slides text and appends them onto a new list
def filter_slide_lists(slides_text, key_word):
    # new_slides_text slices the title page out
    new_slide_list = slides_text[1:]
    filtered_slides_text = []
    for slide in new_slide_list:
        for element in slide:
            if key_word in element.lower() and not any("appendix" in word.lower() for word in slide):
                if key_word == "decision criteria":
                    filtered_slides_text.append(slide[slide.index(element):])
                else:
                    filtered_slides_text.append(slide)
                break
    print(filtered_slides_text)
    return filtered_slides_text

def return_summarized_text(filtered_slides_text):
    summarizer = pipeline('summarization', model = "facebook/bart-large-cnn")
    lessons_learned_text = ""

    if not filtered_slides_text:
        return [{'summary_text': ""}]
    for slide in filtered_slides_text:
        for elements in slide:
            lessons_learned_text += elements

    return summarizer(lessons_learned_text,max_length=150,min_length=60,do_sample=False)

# def create_new_files(slides_text, new_file,summarized_text_file, final_folder_directory):
#     summarized_txt_arr = []
#     lessons_learned_data = filter_slide_lists(slides_text, "lessons learned")
#     summarized_txt_arr.append("Lessons Learned : \n" + return_summarized_text(lessons_learned_data)[0].get("summary_text"))

#     decision_criteria_data = filter_slide_lists(slides_text, "decision criteria")
#     summarized_txt_arr.append("Decision Criteria : \n" + return_summarized_text(decision_criteria_data)[0].get("summary_text"))

#     action_items_data = filter_slide_lists(slides_text, "action item")
#     summarized_txt_arr.append("Action Items : \n" + return_summarized_text(action_items_data)[0].get("summary_text"))

#     extracted_data = decision_criteria_data + lessons_learned_data + action_items_data

#     final_folder_file = os.path.join(final_folder_directory, new_file)
#     final_summarized_file = os.path.join(final_folder_directory, summarized_text_file)

#     fo = open(final_folder_file, "w", encoding='utf-8')
    # for slide in extracted_data:
    #     for list_of_elements in slide:
    #         for elements in list_of_elements:
#                 fo.write(elements)
#             fo.write('\n')
#         fo.write('\n\n')
#     fo.close()

#     sum_txt = open(final_summarized_file,"w", encoding='utf-8')
#     for arr in summarized_txt_arr:
#         sum_txt.write(arr)
#         sum_txt.write('\n\n')
#     fo.close()


def batch_processor(array, folder_path, folder):
    for file in array:
        if file != ".DS_Store":
            summary_file_name = file[:-5]
            summary_file_name += "_summary.docx"
            filename = folder_path + "//" + file
            presentation = Presentation(filename)
            create_new_files(extract_text_from_ppt(presentation), summary_file_name, folder)
            print("Summarization complete")


def create_new_files(slides_text, new_file, final_folder_directory):
    summarized_txt_arr = []
    lessons_learned_data = filter_slide_lists(slides_text, "lessons learned")
    summarized_txt_arr.append("Lessons Learned : \n" + return_summarized_text(lessons_learned_data)[0].get("summary_text"))

    decision_criteria_data = filter_slide_lists(slides_text, "decision criteria")
    summarized_txt_arr.append("Decision Criteria : \n" + return_summarized_text(decision_criteria_data)[0].get("summary_text"))

    action_items_data = filter_slide_lists(slides_text, "action item")
    summarized_txt_arr.append("Action Items : \n" + return_summarized_text(action_items_data)[0].get("summary_text"))

    extracted_data = decision_criteria_data + lessons_learned_data + action_items_data

    final_folder_file = os.path.join(final_folder_directory, new_file)
    document = Document()
    a = document.add_paragraph()
    data = a.add_run("Summarized Version")
    data.bold = True
    for arr in summarized_txt_arr:
        document.add_paragraph(arr)
    b = document.add_paragraph()
    data1 = b.add_run("Scarped Version")
    data1.bold = True
    for slide in extracted_data:
        for list_of_elements in slide:
            document.add_paragraph(list_of_elements)
    document.save(final_folder_file)

